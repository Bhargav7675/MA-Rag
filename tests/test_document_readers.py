"""Tests for corporate document ingestion readers."""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from src.document_readers import (
    IngestReadOptions,
    read_document,
    resolve_pdf_password,
    SUPPORTED_SUFFIXES,
)


def test_supported_suffixes_include_office_formats():
    assert ".docx" in SUPPORTED_SUFFIXES
    assert ".xlsx" in SUPPORTED_SUFFIXES
    assert ".pptx" in SUPPORTED_SUFFIXES


def test_read_txt_and_md(tmp_path: Path):
    text_file = tmp_path / "note.txt"
    text_file.write_text("Hello corporate policy.", encoding="utf-8")
    md_file = tmp_path / "note.md"
    md_file.write_text("# Title\n\nBody text.", encoding="utf-8")

    assert "corporate policy" in read_document(text_file)
    assert "Body text" in read_document(md_file)


def test_read_docx(tmp_path: Path):
    pytest.importorskip("docx")
    from docx import Document

    path = tmp_path / "memo.docx"
    document = Document()
    document.add_paragraph("Quarterly revenue increased.")
    document.save(str(path))

    text = read_document(path)
    assert "Quarterly revenue increased" in text


def test_read_xlsx(tmp_path: Path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    path = tmp_path / "metrics.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Q1"
    sheet["A1"] = "Region"
    sheet["B1"] = "Sales"
    sheet["A2"] = "West"
    sheet["B2"] = 1200000
    workbook.save(path)

    text = read_document(path)
    assert "Q1" in text
    assert "West" in text
    assert "1200000" in text


def test_read_pptx(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(0, 0, 1000000, 500000)
    textbox.text_frame.text = "Launch roadmap for 2026"
    presentation.save(str(path))

    text = read_document(path)
    assert "Launch roadmap for 2026" in text


def test_resolve_pdf_password_prefers_cli(tmp_path: Path):
    passwords_file = tmp_path / "passwords.json"
    passwords_file.write_text(
        json.dumps({"report.pdf": "from-file", "*": "default"}),
        encoding="utf-8",
    )
    options = IngestReadOptions(
        pdf_password="from-cli",
        pdf_passwords_file=passwords_file,
    )
    assert resolve_pdf_password(Path("report.pdf"), options) == "from-cli"


def test_resolve_pdf_password_from_mapping(tmp_path: Path):
    passwords_file = tmp_path / "passwords.json"
    passwords_file.write_text(
        json.dumps({"report.pdf": "secret", "*": "fallback"}),
        encoding="utf-8",
    )
    options = IngestReadOptions(pdf_passwords_file=passwords_file)
    assert resolve_pdf_password(Path("report.pdf"), options) == "secret"
    assert resolve_pdf_password(Path("other.pdf"), options) == "fallback"

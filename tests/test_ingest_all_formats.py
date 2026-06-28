"""Integration: ingest pipeline for every supported document suffix."""

from __future__ import annotations

from pathlib import Path

import pytest

from src.document_readers import IngestReadOptions, read_document
from src.local_retrieval import build_local_index, iter_supported_files


def _write_txt(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def _write_md(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def _write_docx(path: Path, text: str) -> None:
    pytest.importorskip("docx")
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    document.save(str(path))


def _write_xlsx(path: Path, text: str) -> None:
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = text
    workbook.save(path)


def _write_pptx(path: Path, text: str) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(0, 0, 1000000, 500000)
    textbox.text_frame.text = text
    presentation.save(str(path))


def _write_pdf(path: Path, text: str) -> None:
    pytest.importorskip("fitz")
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    document.save(str(path))
    document.close()


def _write_password_pdf(path: Path, text: str, password: str) -> None:
    pytest.importorskip("fitz")
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    document.save(
        str(path),
        encryption=fitz.PDF_ENCRYPT_AES_256,
        user_pw=password,
    )
    document.close()


@pytest.fixture
def format_corpus(tmp_path: Path) -> Path:
    """One file per supported suffix with a unique ingest marker phrase."""
    root = tmp_path / "corpus"
    root.mkdir()
    markers = {
        ".txt": ("notes.txt", "INGEST_MARKER_TXT alpha"),
        ".md": ("notes.md", "# Doc\n\nINGEST_MARKER_MD beta"),
        ".docx": ("memo.docx", "INGEST_MARKER_DOCX gamma"),
        ".xlsx": ("sheet.xlsx", "INGEST_MARKER_XLSX delta"),
        ".pptx": ("deck.pptx", "INGEST_MARKER_PPTX epsilon"),
        ".pdf": ("plain.pdf", "INGEST_MARKER_PDF zeta"),
    }
    writers = {
        ".txt": _write_txt,
        ".md": _write_md,
        ".docx": _write_docx,
        ".xlsx": _write_xlsx,
        ".pptx": _write_pptx,
        ".pdf": _write_pdf,
    }
    for suffix, (name, phrase) in markers.items():
        target = root / name
        writers[suffix](target, phrase)

    secret = root / "secret.pdf"
    _write_password_pdf(secret, "INGEST_MARKER_PDF_SECRET theta", "test-pass-123")
    return root


def test_read_each_supported_suffix(format_corpus: Path) -> None:
    options = IngestReadOptions(pdf_password="test-pass-123", enable_pdf_ocr=False)
    for path in iter_supported_files(format_corpus):
        text = read_document(path, options)
        assert text.strip(), f"No text extracted from {path.name}"


def test_build_local_index_all_formats(format_corpus: Path, tmp_path: Path) -> None:
    index_dir = tmp_path / "index"
    options = IngestReadOptions(pdf_password="test-pass-123", enable_pdf_ocr=False)
    chunk_count = build_local_index(
        format_corpus,
        index_dir=index_dir,
        read_options=options,
    )
    assert chunk_count >= 6

    chunks_file = index_dir / "chunks.jsonl"
    blob = chunks_file.read_text(encoding="utf-8")
    for marker in (
        "INGEST_MARKER_TXT",
        "INGEST_MARKER_MD",
        "INGEST_MARKER_DOCX",
        "INGEST_MARKER_XLSX",
        "INGEST_MARKER_PPTX",
        "INGEST_MARKER_PDF",
        "INGEST_MARKER_PDF_SECRET",
    ):
        assert marker in blob, f"Missing marker {marker} in indexed chunks"

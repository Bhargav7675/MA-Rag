"""Extract plain text from corporate document formats for local ingestion."""

from __future__ import annotations

import io
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from src.env import (
    get_pdf_ocr_enabled,
    get_pdf_ocr_min_chars_per_page,
    get_pdf_password,
    get_pdf_passwords_file,
)

SUPPORTED_SUFFIXES = frozenset(
    {".pdf", ".txt", ".md", ".docx", ".xlsx", ".xls", ".pptx"}
)


@dataclass
class IngestReadOptions:
    """Per-ingest options for password-protected PDFs and OCR."""

    pdf_password: Optional[str] = None
    pdf_passwords_file: Optional[Path] = None
    enable_pdf_ocr: bool = True
    pdf_ocr_min_chars_per_page: int = 50

    @classmethod
    def from_env(
        cls,
        *,
        pdf_password: Optional[str] = None,
        enable_pdf_ocr: Optional[bool] = None,
    ) -> "IngestReadOptions":
        return cls(
            pdf_password=pdf_password or get_pdf_password(),
            pdf_passwords_file=get_pdf_passwords_file(),
            enable_pdf_ocr=(
                get_pdf_ocr_enabled() if enable_pdf_ocr is None else enable_pdf_ocr
            ),
            pdf_ocr_min_chars_per_page=get_pdf_ocr_min_chars_per_page(),
        )


def _load_password_map(path: Path) -> dict[str, str]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError(f"PDF passwords file must be a JSON object: {path}")
    return {str(key): str(value) for key, value in data.items()}


def resolve_pdf_password(
    path: Path,
    options: IngestReadOptions,
) -> Optional[str]:
    if options.pdf_password:
        return options.pdf_password
    if options.pdf_passwords_file and options.pdf_passwords_file.exists():
        mapping = _load_password_map(options.pdf_passwords_file)
        if path.name in mapping:
            return mapping[path.name]
        if "*" in mapping:
            return mapping["*"]
    return get_pdf_password()


def read_document(path: Path, options: Optional[IngestReadOptions] = None) -> str:
    """Return normalized plain text for a supported document path."""
    options = options or IngestReadOptions.from_env()
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        return _read_pdf(path, options)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".xlsx":
        return _read_xlsx(path)
    if suffix == ".xls":
        return _read_xls(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    raise ValueError(f"Unsupported file type: {path}")


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError(
            "DOCX ingestion requires python-docx. Run: pip install -r requirements.txt"
        ) from exc

    document = Document(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError(
            "Excel ingestion requires openpyxl. Run: pip install -r requirements.txt"
        ) from exc

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [
                    str(value).strip()
                    for value in row
                    if value is not None and str(value).strip()
                ]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        workbook.close()
    return "\n\n".join(parts)


def _read_xls(path: Path) -> str:
    try:
        import xlrd
    except ImportError as exc:
        raise RuntimeError(
            "Legacy .xls ingestion requires xlrd. Run: pip install -r requirements.txt"
        ) from exc

    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"## Sheet: {sheet.name}")
        for row_idx in range(sheet.nrows):
            cells = [
                str(sheet.cell_value(row_idx, col_idx)).strip()
                for col_idx in range(sheet.ncols)
                if str(sheet.cell_value(row_idx, col_idx)).strip()
            ]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "PowerPoint ingestion requires python-pptx. "
            "Run: pip install -r requirements.txt"
        ) from exc

    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"## Slide {slide_number}")
            parts.extend(slide_parts)
    return "\n\n".join(parts)


def _read_pdf(path: Path, options: IngestReadOptions) -> str:
    password = resolve_pdf_password(path, options)
    try:
        text, page_count = _extract_pdf_text_pypdf(path, password)
    except Exception:
        text, page_count = _extract_pdf_text_pymupdf(path, password)
    if not options.enable_pdf_ocr:
        return text

    min_chars = options.pdf_ocr_min_chars_per_page
    if page_count == 0:
        return text

    avg_chars = len(text.strip()) / page_count
    if avg_chars >= min_chars:
        return text

    try:
        return _extract_pdf_text_ocr(path, password, min_chars_per_page=min_chars)
    except RuntimeError as exc:
        if text.strip():
            return text
        raise RuntimeError(
            f"PDF {path.name} has little extractable text and OCR failed. "
            f"Install Tesseract (macOS: brew install tesseract). Detail: {exc}"
        ) from exc


def _extract_pdf_text_pypdf(path: Path, password: Optional[str]) -> tuple[str, int]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError(
            "PDF ingestion requires pypdf. Run: pip install -r requirements.txt"
        ) from exc

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        if not password:
            raise ValueError(
                f"PDF is password-protected: {path.name}. "
                "Set MA_RAG_PDF_PASSWORD, MA_RAG_PDF_PASSWORDS_FILE, "
                "or pass --pdf-password to ingest.py."
            )
        status = reader.decrypt(password)
        if status == 0:
            raise ValueError(f"Incorrect password for PDF: {path.name}")

    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages), len(pages)


def _extract_pdf_text_pymupdf(path: Path, password: Optional[str]) -> tuple[str, int]:
    """Fallback PDF text extraction (handles AES-encrypted PDFs via PyMuPDF)."""
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError(
            "PDF ingestion requires pymupdf. Run: pip install -r requirements.txt"
        ) from exc

    document = fitz.open(str(path))
    if document.is_encrypted:
        if not password:
            raise ValueError(
                f"PDF is password-protected: {path.name}. "
                "Set MA_RAG_PDF_PASSWORD, MA_RAG_PDF_PASSWORDS_FILE, "
                "or pass --pdf-password to ingest.py."
            )
        if not document.authenticate(password):
            raise ValueError(f"Incorrect password for PDF: {path.name}")

    pages = [(page.get_text() or "") for page in document]
    page_count = len(pages)
    document.close()
    return "\n\n".join(pages), page_count


def _extract_pdf_text_ocr(
    path: Path,
    password: Optional[str],
    *,
    min_chars_per_page: int,
) -> str:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError(
            "PDF OCR requires pymupdf. Run: pip install -r requirements.txt"
        ) from exc

    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError(
            "PDF OCR requires pytesseract and Pillow. "
            "Run: pip install -r requirements.txt"
        ) from exc

    try:
        pytesseract.get_tesseract_version()
    except pytesseract.TesseractNotFoundError as exc:
        raise RuntimeError(
            "Tesseract OCR binary not found on PATH. "
            "Install it (macOS: brew install tesseract)."
        ) from exc

    document = fitz.open(str(path))
    if document.is_encrypted:
        if not password:
            raise ValueError(
                f"PDF is password-protected: {path.name}. "
                "Set MA_RAG_PDF_PASSWORD or --pdf-password."
            )
        if not document.authenticate(password):
            raise ValueError(f"Incorrect password for PDF: {path.name}")

    parts: list[str] = []
    for page in document:
        text = (page.get_text() or "").strip()
        if len(text) >= min_chars_per_page:
            parts.append(text)
            continue
        pixmap = page.get_pixmap(dpi=200)
        image = Image.open(io.BytesIO(pixmap.tobytes("png")))
        ocr_text = (pytesseract.image_to_string(image) or "").strip()
        parts.append(ocr_text or text)

    document.close()
    return "\n\n".join(parts)
